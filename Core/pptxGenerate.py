import sys
import os
import io
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
import TxtToLog as ttl
import Log
import copy
import six 
from PIL import Image
import imagehash

# 递归获取本页Slide/Shapes 里全部的 text 内容
def checkrecursivelyfortext(shpthissetofshapes,textrun):
    for shape in shpthissetofshapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            textrun=checkrecursivelyfortext(shape.shapes,textrun)
        else:
            if hasattr(shape, "text"):
                textrun.append(shape.text)
    return textrun

# 删除slide （来自https://github.com/scanny/python-pptx/issues/67）
def delete_slide(prs, slide):
    #Make dictionary with necessary information
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

# 替换本页的文本
def replaceTextInShapes(shapes, content, default_replacer, font_change=False, font_name='', font_bold=False, font_size=19):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            replaceTextInShapes(shape.shapes, content, default_replacer, font_change, font_name, font_bold, font_size)
        else:
            if hasattr(shape, 'text'):
                if shape.text.find(default_replacer) != -1:
                    replace_paragraph_text_retaining_initial_formatting(shape.text_frame.paragraphs, content)
                    if font_change:
                        for p in shape.text_frame.paragraphs:
                            p.font.name = font_name
                            p.font.bold = font_bold
                            p.font.size = Pt(font_size)
    return None

# 复制本页至新pres1中
def copy_slide(pres,pres1,index):
    source = pres.slides[index]
    blank_slide_layout = pres.slide_layouts[6]
    dest = pres1.slides.add_slide(blank_slide_layout)
    # print('正在复制模板中第', index, '页面。')

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
        # print('正在复制形状：', shp.name)
        for key, value in six.iteritems(source.part.rels):
                    # Make sure we don't copy a notesSlide relation as that won't exist
                # print('key:   ', key)
                # print('value:   ', value)
                if not "notesSlide" in value.reltype:
                        # if shp.name == '图片 8':
                        #     print('value.reltype: ', value.reltype)
                        #     print('value._target: ', value._target)
                        #     print('value.rId: ', value.rId)
                        #     im = Image.open(io.BytesIO(value._target.image.blob))
                        #     im.show()
                        dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
    return dest

# 测试用- 显示当前页面结构
def showStructure(shapes):
    for idx, shape in enumerate(shapes):
        print('shape中第 ',idx, '个。')
        print('shape: ', shape)
        print('name: ', shape.name)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print('【开始进入Group', shape.name, '】')
            showStructure(shape.shapes)
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            print('【这是图片，图片名为：', shape.image.filename,'】')
            print('cNvPr中id:   ', shape._element._nvXxPr.cNvPr.id)
            print('_pic:   ', shape._pic)
        if hasattr(shape, "text"):
            print('这里包含文字： ', shape.text)
        print('【', shape.name, '结束。】')

# 寻找当页角色立绘对应的shape并返回
def findCharPic(shapes, org_path, bias):
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            result = findCharPic(shape.shapes, org_path, bias)
            if result:
                return result
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            rate = shape.width/shape.height
            img = Image.open(org_path)
            temp = img.resize(shape.image.size)
            current = Image.open(io.BytesIO((shape.image.blob)))
            hash_org = imagehash.average_hash(temp)
            hash_curr = imagehash.average_hash(current)
            dist = abs(hash_curr - hash_org)
            print('新pic为：', shape.image.filename)
            print('dist:', dist)
            print('bias:', bias)
            if dist <= bias:
                return shape
    return None

# ------------------------------------------
# ===  立绘更换测试  ===
# prs = Presentation('恐怖医院模板.pptx')
# bias = 1
# path_org = '林檎普通微笑表情.png'
# path_new = '林檎失控表情.png'
# old_pic = findCharPic(prs.slides[4].shapes, path_org, bias)
# x, y, cx, cy= old_pic.left, old_pic.top, old_pic.width, old_pic.height
# new_pic = prs.slides[4].shapes.add_picture(path_new, x,y, cx, cy)
# old_pic._element.addnext(new_pic._element)
# old_pic._element.getparent().remove(old_pic._element)
# prs.save('result.pptx')

# 立绘替换
def changePic(new_fp, org_fp, slide, bias):
    old_pic = findCharPic(slide.shapes, org_fp, bias)
    x, y, cx, cy= old_pic.left, old_pic.top, old_pic.width, old_pic.height
    new_pic = slide.shapes.add_picture(new_fp, x, y, cx, cy)
    old_pic._element.addnext(new_pic._element)
    old_pic._element.getparent().remove(old_pic._element)
    return slide

# 对应角色名寻找对应的ppt页面
def findModelSlide(slides, name):
    for i, slide in enumerate(slides):
        if is_in_shapes(slide.shapes, name):
            return i
    return 0

# 当前名称是否存在在当前页slide/Shapes中
def is_in_shapes(shapes, name):
    result = False
    texts = checkrecursivelyfortext(shapes,[])
    for text in texts:
        text = text.replace(' ','')
        result = result or text.find(name) != -1 
    return result

# 通过name获取该name对应的slide列表
def getSlideDictionary(slides,names,dictionary):
    for name in names:
        dictionary[name] = findModelSlide(slides, name)
    return dictionary

# 通过character列表获取name列表
def getNamesFromScene(scene):
    if scene.countKP() <= 1:
        return [c.name for c in scene.characters]
    else:
        return [c.PLid if c.is_KP else c.name for c in scene.characters]

# 通过Log生成presentation主函数
def generatePreFromLog(scene, mod_fp, filename, filepath, default_replacer,**args):
    if scene.lines:
        mods = Presentation(mod_fp)
        new_pre = Presentation()
        new_pre.slide_width = mods.slide_width
        new_pre.slide_height = mods.slide_height
        dictionary = getSlideDictionary(mods.slides, getNamesFromScene(scene), {})
        i = 0
        for l in scene.lines:
            name = l.character.PLid if l.character.is_KP and scene.countKP()>1 else l.character.name
            index = dictionary[name]
            replaceTextInShapes(mods.slides[index].shapes, l.content, default_replacer, **args)
            copy_slide(mods, new_pre, index)
            if l.pic_fp:
                print(l.content)
                print(index)
                old_fp = l.character.imges[l.character.original_img]
                changePic(l.pic_fp, old_fp, new_pre.slides[-1], 1)
            mods = Presentation(mod_fp)
        fp = os.path.join(filepath, filename + '.pptx')
        new_pre.save(fp)
    else:
        print("台词列表为空。") 

# 官方Issues: https://github.com/scanny/python-pptx/issues/285
# 不更改字体的前提下改变文字内容
def replace_paragraph_text_retaining_initial_formatting(paragraphs, new_text):
    if paragraphs:
        paragraph = paragraphs[0]
        for p in paragraphs[1:]:
            p.clear()   
    else:
        print('空paragraphs列表。')
        return None
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text

# prs = Presentation('model.pptx')

# pre = ''
# for s in prs.slides:
#     for shape in s.shapes:
#         if hasattr(shape, "text"):
#             a = shape.text
#             shape.text = 'what???'
#             a = 'abcde'
#             print('yes')

# d_replacer = '這一周也不例外，你们刚刚从家里，神社，公园来到集会地点，现在距离集会开始還有一段时間，你们可以自由交流一下'
# kp = Log.KP('KP','')
# zhima = Log.PL('芝麻','', '阿恬')
# qiuqiu = Log.PL('球球','', '蓝蓝')
# momo = Log.PL('momo','','Syin')
# touzi = Log.Dice('骰子')

# path_characters = 'D:/TRPG/TRPG_replay_video_generator/charlist.json'
# ttl.saveCharacterList([kp, zhima, qiuqiu, momo, touzi], path_characters)
# charlist = ttl.getCharacterList(path_characters)
# apple =  Log.PL('林檎', '', 'EE')
# apple.set_imgpath('D:/TRPG/TRPG_replay_video_generator/林檎')
# paths = apple.imges
# d = {}
# for path in paths:
#     if path.find('失控') != -1:
#         d[path] = ['愤怒的']
# apple.set_dictionary(d)
# apple.original_img = 2
# charlist = [Log.KP('火火', ''), apple]
# txt_fp = 'D:/TRPG/TRPG_replay_video_generator/医院05.txt'
# txt_format = 'ldn'
# txt_encoding = 'GBK'
# lines = ttl.readFromTxt(txt_fp, charlist,path_characters, txt_format=txt_format, encoding=txt_encoding)
# scene = Log.Scene(lines, characters=charlist)
# model_ppt_fp = '恐怖医院模板.pptx'
# prs = Presentation(model_ppt_fp)
# dictionary = getSlideDictionary(prs.slides, [c.name for c in charlist], {})
# name = '123'
# fp = 'C:/Users/44418/Desktop'
# font_change = True
# font_size = 20
# font_bold = True
# font_name = '清松手寫體1'
# d_replacer = '我有自行车'
# generatePreFromLog(scene,model_ppt_fp,name, fp, d_replacer)


# for slide in prs.slides:
#     print(pre + ' 【start slide】\n')
#     print(checkrecursivelyfortext(slide.shapes,[], pre))
#     print(pre + '【end slide】')
#     print(' ')